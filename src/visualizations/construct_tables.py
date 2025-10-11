# -*- coding: utf-8 -*-
"""
Export regression / Granger-causality results to LaTeX or HTML.

Key features
------------
- output_format: 'tex' (LaTeX) or 'html'
- tex_mode: 'longtable' (standalone, non-floating) or 'tabular' (to \input inside a floating table)
- Clean LaTeX output for β/γ headers (no Unicode errors)
- NaN rendered as blanks
- Text cells safely escaped (so escape=False is usable to allow math in headers)

Also includes simple PPTX helpers (optional).
"""
from __future__ import annotations

from pathlib import Path
from typing import Sequence
import pandas as pd

# Optional PPTX helpers
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ----------------------------- configuration -------------------------------- #
_OUTPUT_ORDER = [
    "Model", "AINI_variant", "Period", "Ticker",
    "β₀", "β1", "β2", "β3", "γ1", "γ2", "γ3",
    "R^2", "R^2_adj",
    "Analytical F_stat p", "Empirical F_stat p",
    "BH Analytical p", "BH Empirical p",
]

_DROP_COLS = [
    "Direction", "Lags",
    "N_obs", "N_boot", "N_boot_valid",
    "F_stat", "df_num", "df_den",
    "BH_reject_F", "BH_reject_F_HC3",
    "joint rej. (α=0.1)",
]

# Columns that should be right-aligned (numeric-like)
_NUMERIC_LIKE = {
    "β₀", "β1", "β2", "β3", "γ1", "γ2", "γ3",
    "R^2", "R^2_adj",
    "Analytical F_stat p", "Empirical F_stat p",
    "BH Analytical p", "BH Empirical p",
}

# ------------------------------- helpers ------------------------------------ #
def _project_tables_dir() -> Path:
    here = Path(__file__).resolve()
    project_root = here.parents[2]  # .../AI_narrative_index
    return project_root / "reports" / "tables"


def _stars_from_p(p: float) -> str:
    """Return significance stars for p-value."""
    try:
        if pd.isna(p):
            return ""
        p = float(p)
    except Exception:
        return ""
    if p < 0.01:
        return "***"
    if p < 0.05:
        return "**"
    if p < 0.10:
        return "*"
    return ""


def _fmt_p(p: float, digits: int = 3) -> str:
    """Format p-value with fixed decimals and floor at <0.00..01."""
    try:
        if pd.isna(p):
            return ""
        p = float(p)
        floor = 10 ** (-digits)
        return f"<0.{('0'*(digits-1))}1" if p < floor else f"{p:.{digits}f}"
    except Exception:
        return ""


def _column_format(columns: Sequence[str]) -> str:
    """
    Build LaTeX column spec: right-align numeric-like, left-align everything else.
    Uses the *logical* column names (before math sanitization).
    """
    spec = ["r" if c in _NUMERIC_LIKE else "l" for c in columns]
    return "@{}" + "".join(spec) + "@{}"


def _escape_latex_text(s: str) -> str:
    """
    Escape LaTeX specials in text cells so we can set escape=False overall
    (needed to allow math in headers). NaN/None -> "".
    """
    if s is None or pd.isna(s):
        return ""
    s = str(s)
    return (s.replace("\\", r"\textbackslash{}")
             .replace("&", r"\&")
             .replace("%", r"\%")
             .replace("$", r"\$")
             .replace("#", r"\#")
             .replace("_", r"\_")
             .replace("{", r"\{")
             .replace("}", r"\}")
             .replace("~", r"\textasciitilde{}")
             .replace("^", r"\textasciicircum{}")
             .replace("<", r"\textless{}")
             .replace(">", r"\textgreater{}"))


def _sanitize_latex_header(col: str) -> str:
    """
    Map Unicode β/γ + subscripts to LaTeX math headers.
    Only column headers are transformed; body stays as text (escaped).
    """
    mapping = {
        "β₀": r"$\beta_0$", "β1": r"$\beta_1$", "β2": r"$\beta_2$", "β3": r"$\beta_3$",
        "γ1": r"$\gamma_1$", "γ2": r"$\gamma_2$", "γ3": r"$\gamma_3$",
    }
    return mapping.get(col, col)

# ------------------------------- main API ----------------------------------- #
def export_regression_table(
    df: pd.DataFrame,
    title: str,
    output_filename: str,
    *,
    p_digits: int = 3,
    coef_digits: int | None = None,
    font_size_cmd: str = "tiny",
    tabcolsep_pt: float = 2.0,
    sort_by: Sequence[str] | None = None,
    output_format: str = "tex",          # 'tex' or 'html'
    tex_mode: str = "longtable",         # 'longtable' or 'tabular'
    tex_include_caption: bool = True,    # only used when tex_mode='longtable'
) -> Path:
    """
    Export regression results to LaTeX or HTML.

    - For a floating LaTeX table:
        call with tex_mode='tabular', tex_include_caption=False,
        then wrap the generated .tex file:
            \begin{table}[ht]
              \centering
              \resizebox{\textwidth}{!}{\input{.../file.tex}}
              \caption{... S\&P~500 ...}
              \label{...}
            \end{table}

    - For a standalone longtable:
        call with tex_mode='longtable' (default) and let this function add caption/label.
        Do NOT wrap it in a floating table.
    """
    if not isinstance(df, pd.DataFrame):
        raise TypeError("df must be a pandas.DataFrame")
    if output_format not in {"tex", "html"}:
        raise ValueError("output_format must be 'tex' or 'html'")
    if output_format == "tex" and tex_mode not in {"longtable", "tabular"}:
        raise ValueError("tex_mode must be 'longtable' or 'tabular'")

    outdir = _project_tables_dir()
    outdir.mkdir(parents=True, exist_ok=True)

    if not output_filename.endswith(f".{output_format}"):
        output_filename = f"{output_filename}.{output_format}"
    out_path = outdir / output_filename

    # ---------------- normalize schema & values ---------------- #
    work = df.copy()

    rename_map = {
        "r2_u": "R^2",
        "adj_r2_u": "R^2_adj",
        "Original_F_pval": "Analytical F_stat p",
        "Empirical_F_pval": "Empirical F_stat p",
        "BH analytical p": "BH Analytical p",
        "BH empirical p": "BH Empirical p",
    }
    for src, dst in rename_map.items():
        if src in work.columns:
            work[dst] = work[src]

    # coefficient formatting
    if coef_digits is not None:
        for c in ["β₀", "β1", "β2", "β3", "γ1", "γ2", "γ3"]:
            if c in work.columns:
                work[c] = work[c].apply(lambda v: "" if pd.isna(v) else f"{float(v):.{coef_digits}f}")

    # plain p-values
    for pcol in ["Analytical F_stat p", "Empirical F_stat p"]:
        if pcol in work.columns:
            work[pcol] = work[pcol].apply(lambda v: _fmt_p(v, p_digits))

    # BH p-values + stars based on BH Empirical p
    stars_source = None
    if "BH Empirical p" in work.columns:
        stars_source = work["BH Empirical p"].apply(_stars_from_p)
        work["BH Empirical p"] = work["BH Empirical p"].apply(lambda v: _fmt_p(v, p_digits))
    if "BH Analytical p" in work.columns:
        work["BH Analytical p"] = work["BH Analytical p"].apply(lambda v: _fmt_p(v, p_digits))
    if stars_source is not None:
        if "BH Empirical p" in work.columns:
            work["BH Empirical p"] = work["BH Empirical p"] + stars_source
        if "BH Analytical p" in work.columns:
            work["BH Analytical p"] = work["BH Analytical p"] + stars_source

    # drop unnecessary cols
    work = work.drop(columns=[c for c in _DROP_COLS if c in work.columns], errors="ignore")

    # column order
    cols = [c for c in _OUTPUT_ORDER if c in work.columns]
    work = work[cols]

    # optional sorting
    if sort_by:
        missing = [c for c in sort_by if c not in work.columns]
        if not missing:
            work = work.sort_values(list(sort_by)).reset_index(drop=True)

    # ----------------------------- EXPORT ----------------------------- #
    if output_format == "tex":
        # Build column spec BEFORE header sanitization (based on logical names)
        colspec = _column_format(work.columns)

        # Sanitize headers to LaTeX math (β/γ), avoid Unicode subscripts
        header_cols = [_sanitize_latex_header(c) for c in work.columns]

        # Sanitize text cells so we can set escape=False; blank out NaNs
        obj_cols = list(work.select_dtypes(include="object").columns)
        for c in obj_cols:
            work[c] = work[c].apply(_escape_latex_text)

        # Replace DataFrame headers with sanitized versions
        work.columns = header_cols

        # Prepare to_latex kwargs
        to_latex_kwargs = dict(
            index=False,
            escape=False,       # allow math in headers
            na_rep="",          # blank instead of NaN
            column_format=colspec,
            bold_rows=False,
            multicolumn=False,
            multicolumn_format="c",
            longtable=(tex_mode == "longtable"),
        )
        if tex_mode == "longtable" and tex_include_caption:
            to_latex_kwargs["caption"] = title
            to_latex_kwargs["label"] = f"tab:{out_path.stem}"

        latex_core = work.to_latex(**to_latex_kwargs)

        if tex_mode == "longtable":
            # A4-friendly tweaks only for longtable
            preamble = (
                "% AUTO-GENERATED by export_regression_table\n"
                "\\begingroup\n"
                "\\setlength\\LTleft{0pt}\\setlength\\LTright{0pt}\n"
                f"\\setlength\\tabcolsep{{{tabcolsep_pt}pt}}\n"
                f"\\{font_size_cmd}\n"
            )
            postamble = "\n\\endgroup\n"
            out_path.write_text(preamble + latex_core + postamble, encoding="utf-8")
        else:
            # Plain tabular — safe to wrap in \begin{table} ... \end{table}
            header = "% AUTO-GENERATED TABULAR (no caption/label) — safe to wrap inside a floating table\n"
            out_path.write_text(header + latex_core, encoding="utf-8")

    elif output_format == "html":
        html_core = (
            work.style
            .set_caption(title)
            .set_table_styles([
                {"selector": "caption", "props": [("caption-side", "top"), ("font-weight", "bold"), ("font-size", "120%")]},
                {"selector": "th, td", "props": [("border", "1px solid #ddd"), ("padding", "4px 6px"), ("text-align", "center")]},
                {"selector": "thead", "props": [("background-color", "#f2f2f2")]}
            ])
            .hide(axis="index")
            .to_html()
        )
        html_doc = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8"><title>{title}</title>
<style>
body {{ font-family: Arial, sans-serif; margin: 40px; }}
table {{ border-collapse: collapse; width: 100%; font-size: 12px; }}
</style>
</head>
<body>
<h2>{title}</h2>
{html_core}
</body>
</html>
"""
        out_path.write_text(html_doc, encoding="utf-8")

    return out_path

# ---------------- PPTX helpers (optional) ---------------- #
BOLD_IF_POSITIVE_COLS = {'β₁','β₂','β₃','γ₁','γ₂','γ₃','β1','β2','β3','γ1','γ2','γ3'}

def add_text_to_cell(cell, text: str, bold=False, base_size=12):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(base_size)
    run.font.bold = bold

def add_table_slide(prs: Presentation, title: str, headers, rows, font_size=12):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    p = txbox.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True

    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.4 + 0.3*n_rows)
    ).table

    for j, h in enumerate(headers):
        add_text_to_cell(tbl.cell(0, j), str(h), bold=False, base_size=font_size)
        for run in tbl.cell(0, j).text_frame.paragraphs[0].runs:
            run.font.bold = True

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            col_name = headers[j]
            text = "not tested" if pd.isna(val) else str(val)
            make_bold = False
            if col_name in BOLD_IF_POSITIVE_COLS and pd.notna(val):
                try:
                    make_bold = float(val) > 0
                except Exception:
                    make_bold = False
            add_text_to_cell(tbl.cell(i, j), text, bold=make_bold, base_size=font_size)

    for j in range(n_cols):
        tbl.columns[j].width = Inches(9.0 / n_cols)
    return slide

def df_to_pptx(df: pd.DataFrame, outpath="df_tables.pptx", rows_per_slide=12):
    prs = Presentation()
    headers = list(df.columns)
    for start in range(0, len(df), rows_per_slide):
        chunk = df.iloc[start:start+rows_per_slide]
        rows = chunk.fillna("not tested").values.tolist()
        add_table_slide(prs, f"Rows {start+1}-{start+len(chunk)}", headers, rows, font_size=12)
    prs.save(outpath)
    return Path(outpath).resolve()
