# -*- coding: utf-8 -*-
import re
from pathlib import Path

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# ========= 1) REJECTION RATE LOGIC =========

def _to_bool_series(s: pd.Series) -> pd.Series:
    """Robustly coerce to boolean."""
    if pd.api.types.is_bool_dtype(s):
        return s
    if pd.api.types.is_numeric_dtype(s):
        return s.fillna(0).astype(float) > 0
    s_str = s.astype(str).str.strip().str.lower()
    return s_str.isin({"true", "1", "yes", "y", "t"})

def _agg_rates(df: pd.DataFrame, flag_col: str, key: str) -> pd.DataFrame:
    """sum, count, rate[%] by key."""
    out = (
        df.groupby(key)[flag_col]
          .agg(sum="sum", count="count")
          .assign(rate=lambda x: 100 * x["sum"] / x["count"])
          .reset_index()
          .sort_values("rate", ascending=False, key=lambda s: s.astype(float))
    )
    return out

def compute_rejection_rates_all(gc_all_results: pd.DataFrame, direction_substr="AINI_to_RET"):
    """
    Returns:
      {
        'by_year':   {'boot': df, 'hc3': df, 'both': df},
        'by_ticker': {'boot': df, 'hc3': df, 'both': df},
        'a2r': filtered df with flags
      }
    """
    a2r = gc_all_results[gc_all_results["Direction"].str.contains(direction_substr, case=False)].copy()
    a2r["Year"]   = a2r["Year"].astype(str)
    a2r["Ticker"] = a2r["Ticker"].astype(str)

    a2r["rej_bh_boot"] = _to_bool_series(a2r["BH_reject_F"])
    a2r["rej_bh_hc"]   = _to_bool_series(a2r["BH_reject_F_HC3"])
    a2r["rej_both"]    = a2r["rej_bh_boot"] & a2r["rej_bh_hc"]

    by_year = {
        "boot": _agg_rates(a2r, "rej_bh_boot", "Year"),
        "hc3":  _agg_rates(a2r, "rej_bh_hc",   "Year"),
        "both": _agg_rates(a2r, "rej_both",    "Year"),
    }
    by_ticker = {
        "boot": _agg_rates(a2r, "rej_bh_boot", "Ticker"),
        "hc3":  _agg_rates(a2r, "rej_bh_hc",   "Ticker"),
        "both": _agg_rates(a2r, "rej_both",    "Ticker"),
    }
    return {"by_year": by_year, "by_ticker": by_ticker, "a2r": a2r}


# ========= 2) ADD “TRADING DAYS” COLUMN =========

# Your mapping (use underscores like in your Year labels)
TRADING_DAYS_MAP = {
    "2024_25":      365,   # "2024-25"
    "2023_24_25":   553,   # "2023-2025"
    "2025":         113,
    "2023":         188,
    "2023_24":      440,   # "2023-2024"
    "2024":         252,
}

def add_trading_days_columns(rates: dict) -> dict:
    """
    Adds 'Trading days' column to:
      - by_year tables: via Year→days map
      - by_ticker tables: sum of days across UNIQUE Year-periods present for each ticker
    Returns a NEW dict (does not mutate input).
    """
    a2r = rates["a2r"].copy()

    # ---- by_year: direct map
    out_by_year = {}
    for method, df in rates["by_year"].items():
        df2 = df.copy()
        df2["Trading days"] = df2["Year"].map(TRADING_DAYS_MAP).fillna(0).astype(int)
        # optional: order columns
        cols = ["Year", "sum", "count", "rate", "Trading days"]
        out_by_year[method] = df2[cols]

    # ---- by_ticker: sum trading days across UNIQUE Year-periods for each ticker
    # unique (Ticker, Year)
    uniq = a2r[["Ticker", "Year"]].drop_duplicates()
    uniq["days"] = uniq["Year"].map(TRADING_DAYS_MAP).fillna(0).astype(int)
    ticker_days = uniq.groupby("Ticker", as_index=False)["days"].sum().rename(columns={"days": "Trading days"})

    out_by_ticker = {}
    for method, df in rates["by_ticker"].items():
        df2 = df.copy()
        df2 = df2.merge(ticker_days, on="Ticker", how="left")
        df2["Trading days"] = df2["Trading days"].fillna(0).astype(int)
        cols = ["Ticker", "sum", "count", "rate", "Trading days"]
        out_by_ticker[method] = df2[cols].sort_values("rate", ascending=False, key=lambda s: s.astype(float))

    return {"by_year": out_by_year, "by_ticker": out_by_ticker, "a2r": a2r}


# ========= 3) POWERPOINT (TABLES ONLY, NO PLOTS) =========

def _add_title(slide, text, fontsize=24):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(fontsize)
        r.font.bold = True

def _add_table(slide, df: pd.DataFrame, top_left=(0.5, 1.2), col_width_total=9.0, row_height=0.35, font_size=12):
    """
    Insert a PPT table from DataFrame. Formats 'rate' as xx.xx%.
    """
    n_rows = len(df) + 1
    n_cols = len(df.columns)
    left, top = Inches(top_left[0]), Inches(top_left[1])
    width, height = Inches(col_width_total), Inches(0.5 + row_height * n_rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # column widths
    for j in range(n_cols):
        table.columns[j].width = Inches(col_width_total / n_cols)

    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        p = cell.text_frame.paragraphs[0]
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(font_size)

    # body
    for i in range(len(df)):
        for j, col in enumerate(df.columns):
            val = df.iloc[i, j]
            if col == "rate":
                txt = f"{float(val):.2f}%"
            elif pd.api.types.is_number(val):
                txt = f"{val:.0f}" if float(val).is_integer() else f"{val:.4f}"
            else:
                txt = str(val)
            cell = table.cell(i + 1, j)
            cell.text = txt
            p = cell.text_frame.paragraphs[0]
            for r in p.runs:
                r.font.size = Pt(font_size)

def _add_table_slide(prs: Presentation, title: str, df: pd.DataFrame):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_title(slide, title)
    _add_table(slide, df, top_left=(0.5, 1.2))


def export_rejection_rates_to_pptx_all_tables_only(
    rates_with_days: dict,
    outpath="rejection_rates_all_methods_tables_only.pptx",
    top_tickers=15
):
    """
    Creates slides with tables only (no plots), and includes 'Trading days' column.
    """
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "AINI → Returns: Rejection Rates (Bootstrap, HC3, Both) — Tables Only")

    # By Year
    by_year = rates_with_days["by_year"]
    _add_table_slide(prs, "By Year — Bootstrap (BH on F_boot)", by_year["boot"])
    _add_table_slide(prs, "By Year — HC3 (BH on F_analytic)",   by_year["hc3"])
    _add_table_slide(prs, "By Year — Both (intersection)",      by_year["both"])

    # By Ticker (Top N)
    by_ticker = rates_with_days["by_ticker"]
    _add_table_slide(prs, f"By Ticker — Bootstrap (Top {top_tickers})",
                     by_ticker["boot"].sort_values("rate", ascending=False).head(top_tickers).reset_index(drop=True))
    _add_table_slide(prs, f"By Ticker — HC3 (Top {top_tickers})",
                     by_ticker["hc3"].sort_values("rate", ascending=False).head(top_tickers).reset_index(drop=True))
    _add_table_slide(prs, f"By Ticker — Both (Top {top_tickers})",
                     by_ticker["both"].sort_values("rate", ascending=False).head(top_tickers).reset_index(drop=True))

    outpath = Path(outpath)
    outpath.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(outpath))
    return outpath.resolve()


# ========= 4) USAGE =========
# from this_module import compute_rejection_rates_all, add_trading_days_columns, export_rejection_rates_to_pptx_all_tables_only
#
# rates = compute_rejection_rates_all(gc_all_results, direction_substr="AINI_to_RET")
# rates_with_days = add_trading_days_columns(rates)
# ppt = export_rejection_rates_to_pptx_all_tables_only(
#     rates_with_days,
#     outpath="reports/figures/rejection_rates_all_methods_tables_only.pptx",
#     top_tickers=15
# )
# print("Saved PPT:", ppt)
