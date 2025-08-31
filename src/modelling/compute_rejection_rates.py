# -*- coding: utf-8 -*-
import re
from pathlib import Path

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ========== 1) COMPUTE REJECTION RATES ==========

def _to_bool_series(s: pd.Series) -> pd.Series:
    """
    Robustly convert a column to boolean:
    - already bool: return as-is
    - numeric 0/1 or floats: >0 => True
    - strings like 'True'/'False'/'1'/'0': map accordingly
    Otherwise, coerce to False.
    """
    if pd.api.types.is_bool_dtype(s):
        return s
    if pd.api.types.is_numeric_dtype(s):
        return s.fillna(0).astype(float) > 0
    # strings
    s_str = s.astype(str).str.strip().str.lower()
    return s_str.isin({"true", "1", "yes", "y", "t"})

def compute_rejection_rates(gc_all_results: pd.DataFrame, direction_substr="AINI_to_RET"):
    """
    Filters to AINI->RET rows and computes rejection rates:
    - rej_bh_boot := BH_reject_F  (bool)
    - rej_bh_hc   := BH_reject_F_HC3 (bool)
    - rej_both    := rej_bh_boot & rej_bh_hc

    Returns:
        by_year_both, by_ticker_both  (two DataFrames with sum, count, rate [%])
    """
    a2r = gc_all_results[gc_all_results["Direction"].str.contains(direction_substr, case=False)].copy()

    # normalize types
    a2r["Year"] = a2r["Year"].astype(str)
    a2r["Ticker"] = a2r["Ticker"].astype(str)

    a2r["rej_bh_boot"] = _to_bool_series(a2r["BH_reject_F"])
    a2r["rej_bh_hc"] = _to_bool_series(a2r["BH_reject_F_HC3"])
    a2r["rej_both"] = a2r["rej_bh_boot"] & a2r["rej_bh_hc"]

    def _agg_rates(df: pd.DataFrame, key: str) -> pd.DataFrame:
        out = (
            df.groupby(key)["rej_both"]
              .agg(sum="sum", count="count")
              .assign(rate=lambda x: 100 * x["sum"] / x["count"])
              .sort_values("rate", ascending=False)
              .reset_index()
        )
        return out

    by_year = _agg_rates(a2r, "Year")
    by_ticker = _agg_rates(a2r, "Ticker")
    return by_year, by_ticker


# ========== 2) PLOTTING HELPERS (optional bar charts) ==========

def save_bar_chart(df: pd.DataFrame, x_col: str, y_col: str, outpath: Path, title: str = "", xlabel: str = "", ylabel: str = ""):
    """
    Save a simple bar chart as PNG from df[x_col] vs df[y_col].
    Note: Keep a single axes, no style or custom colors (clean and portable).
    """
    fig, ax = plt.subplots(figsize=(10, 4))
    ax.bar(df[x_col].astype(str), df[y_col])
    ax.set_title(title)
    ax.set_xlabel(x_col if not xlabel else xlabel)
    ax.set_ylabel(ylabel if ylabel else y_col)
    ax.set_ylim(0, max(100, df[y_col].max() * 1.1) if len(df) else 100)
    ax.tick_params(axis='x', rotation=45)
    fig.tight_layout()
    outpath.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(outpath, dpi=200, bbox_inches="tight")
    plt.close(fig)


# ========== 3) POWERPOINT HELPERS ==========

def _add_title(slide, text, fontsize=24):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(fontsize)
        r.font.bold = True

def _add_table(slide, df: pd.DataFrame, top_left=(0.5, 1.2), col_width_total=9.0, row_height=0.35, font_size=12):
    """
    Insert a PPT table from DataFrame.
    """
    n_rows = len(df) + 1  # header + rows
    n_cols = len(df.columns)

    left, top = Inches(top_left[0]), Inches(top_left[1])
    width, height = Inches(col_width_total), Inches(0.5 + row_height * n_rows)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # set uniform column widths
    for j in range(n_cols):
        table.columns[j].width = Inches(col_width_total / n_cols)

    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        p = cell.text_frame.paragraphs[0]
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(font_size)

    # body
    for i in range(len(df)):
        for j, col in enumerate(df.columns):
            val = df.iloc[i, j]
            # format floats: rate with 2 decimals and a % sign if the column is 'rate'
            if col == "rate":
                txt = f"{float(val):.2f}%"
            elif pd.api.types.is_number(val):
                txt = f"{val:.0f}" if float(val).is_integer() else f"{val:.4f}"
            else:
                txt = str(val)
            cell = table.cell(i + 1, j)
            cell.text = txt
            p = cell.text_frame.paragraphs[0]
            for r in p.runs:
                r.font.size = Pt(font_size)

def _add_image(slide, img_path: Path, left=0.5, top=4.5, width=9.0):
    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(width))

def export_rejection_rates_to_pptx(
    by_year: pd.DataFrame,
    by_ticker: pd.DataFrame,
    outpath="rejection_rates.pptx",
    top_tickers=15,
    make_charts=True,
):
    prs = Presentation()

    # --- Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _add_title(slide, "AINI → Returns: Rejection Rates (both-method)")

    # --- By Year
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "Rejection Rate by Year (both-method)")
    _add_table(slide, by_year, top_left=(0.5, 1.2))

    # optional bar chart
    if make_charts and not by_year.empty:
        img_year = Path("reports/figures/rej_rate_by_year.png")
        save_bar_chart(by_year, x_col="Year", y_col="rate", outpath=img_year, title="Rejection Rate by Year", ylabel="Rate (%)")
        _add_image(slide, img_year, left=0.5, top=4.0, width=9.0)

    # --- By Ticker (top N)
    topN = by_ticker.sort_values("rate", ascending=False).head(top_tickers).reset_index(drop=True)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, f"Rejection Rate by Ticker (Top {top_tickers}, both-method)")
    _add_table(slide, topN, top_left=(0.5, 1.2))

    if make_charts and not topN.empty:
        img_tick = Path("reports/figures/rej_rate_by_ticker_top.png")
        save_bar_chart(topN, x_col="Ticker", y_col="rate", outpath=img_tick, title=f"Top {top_tickers} Tickers by Rejection Rate", ylabel="Rate (%)")
        _add_image(slide, img_tick, left=0.5, top=4.0, width=9.0)

    outpath = Path(outpath)
    outpath.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(outpath))
    return outpath.resolve()


# ========== 4) END-TO-END EXAMPLE ==========

if __name__ == "__main__":
    # Replace this with your actual DataFrame load
    # e.g., gc_all_results = pd.read_csv("path/to/gc_all_results.csv")
    raise SystemExit("Import this module and call compute_rejection_rates(...) / export_rejection_rates_to_pptx(...).")
